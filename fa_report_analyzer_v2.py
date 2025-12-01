"""
FA Report Analyzer v2.0 - 失效分析報告評估工具
支援地端 Ollama 模型、OpenAI API 和圖片解析功能
"""

import json
import base64
import anthropic
from pathlib import Path
from datetime import datetime
from typing import Dict, List, Tuple, Optional, Any
import sys
import io

try:
    import pandas as pd
except ImportError:
    print("需要安裝 pandas: pip install pandas --break-system-packages")
    sys.exit(1)

# 可選依賴
try:
    from PIL import Image
    HAS_PIL = True
except ImportError:
    HAS_PIL = False

try:
    import ollama
    HAS_OLLAMA = True
except ImportError:
    HAS_OLLAMA = False

try:
    from openai import OpenAI
    HAS_OPENAI = True
except ImportError:
    HAS_OPENAI = False


class FAReportAnalyzer:
    """FA 報告分析器 v2.0 - 支援多種 LLM 後端和圖片解析"""
    
    def __init__(self,
                 backend: str = "ollama",
                 model: str = None,
                 api_key: str = None,
                 base_url: str = None,
                 skip_images: bool = False):
        """初始化分析器

        Args:
            backend: LLM 後端 ('ollama', 'openai', 'anthropic')
            model: 模型名稱
            api_key: API key (OpenAI/Anthropic 需要)
            base_url: API base URL (OpenAI 相容接口)
            skip_images: 是否跳過圖片分析 (僅分析文字)
        """
        self.backend = backend.lower()
        self.api_key = api_key
        self.base_url = base_url
        self.skip_images = skip_images
        self.temp_files = []  # 用於追蹤需要清理的臨時文件
        
        # 設定預設模型
        if model:
            self.model = model
        else:
            if self.backend == "ollama":
                # self.model = "gpt-oss:20b"  # 支援視覺的模型
                self.model = "llama3.1:latest"
            elif self.backend == "openai":
                # self.model = "gpt-4.1-mini“
                # self.model = "gpt-4o-2024-05-13"
                self.model = "gpt-4o-mini-2024-07-18"
            elif self.backend == "anthropic":
                self.model = "claude-sonnet-4-20250514"
            else:
                self.model = "llama3.2-vision:latest"
        
        # 初始化客戶端
        self._init_client()
        
        # 評估維度與權重
        self.dimensions = {
            "基本資訊完整性": 15,
            "問題描述與定義": 15,
            "分析方法與流程": 20,
            "數據與證據支持": 20,
            "根因分析": 20,
            "改善對策": 10
        }
        
        # 評分標準
        self.grade_criteria = {
            'A': (90, 100, '卓越報告'),
            'B': (80, 89, '良好報告'),
            'C': (70, 79, '合格報告'),
            'D': (60, 69, '待改進報告'),
            'F': (0, 59, '不合格報告')
        }
    
    def _init_client(self):
        """初始化 LLM 客戶端"""
        if self.backend == "ollama":
            if not HAS_OLLAMA:
                print("警告: 需要安裝 ollama: pip install ollama --break-system-packages")
                print("嘗試使用其他後端...")
                self.backend = "anthropic" if self.api_key else None
                if not self.backend:
                    raise RuntimeError("無可用的 LLM 後端")
            else:
                self.client = None  # Ollama 使用函數式調用
                print(f"✓ 使用 Ollama 地端模型: {self.model}")
        
        elif self.backend == "openai":
            if not HAS_OPENAI:
                raise ImportError("需要安裝 openai: pip install openai --break-system-packages")
            self.client = OpenAI(
                api_key=self.api_key,
                base_url=self.base_url
            )
            print(f"✓ 使用 OpenAI API: {self.model}")
        
        elif self.backend == "anthropic":
            self.client = anthropic.Anthropic(api_key=self.api_key)
            print(f"✓ 使用 Anthropic Claude: {self.model}")
        
        else:
            raise ValueError(f"不支援的後端: {self.backend}")
    
    def _convert_ppt_to_pptx(self, ppt_path: str) -> Optional[str]:
        """嘗試將 .ppt 轉換為 .pptx
        
        Args:
            ppt_path: .ppt 文件路徑
            
        Returns:
            轉換後的 .pptx 文件路徑，失敗返回 None
        """
        import subprocess
        import os
        
        pptx_path = ppt_path.rsplit('.', 1)[0] + '_converted.pptx'
        
        # 方法 1: 嘗試使用 LibreOffice
        try:
            # 檢查 LibreOffice 是否安裝
            libreoffice_paths = [
                'libreoffice',
                '/Applications/LibreOffice.app/Contents/MacOS/soffice',  # macOS
                '/usr/bin/libreoffice',  # Linux
                'C:\\Program Files\\LibreOffice\\program\\soffice.exe',  # Windows
            ]
            
            libreoffice_cmd = None
            for path in libreoffice_paths:
                try:
                    if os.path.exists(path) or subprocess.run(
                        [path, '--version'], 
                        capture_output=True, 
                        timeout=2
                    ).returncode == 0:
                        libreoffice_cmd = path
                        break
                except:
                    continue
            
            if libreoffice_cmd:
                print(f"  使用 LibreOffice 進行轉換...")
                output_dir = os.path.dirname(ppt_path)
                result = subprocess.run(
                    [
                        libreoffice_cmd,
                        '--headless',
                        '--convert-to', 'pptx',
                        '--outdir', output_dir,
                        ppt_path
                    ],
                    capture_output=True,
                    timeout=30
                )
                
                if result.returncode == 0:
                    # LibreOffice 會生成與原文件同名但副檔名為 .pptx 的文件
                    auto_pptx = ppt_path.rsplit('.', 1)[0] + '.pptx'
                    if os.path.exists(auto_pptx):
                        self.temp_files.append(auto_pptx)  # 記錄臨時文件
                        return auto_pptx
                    elif os.path.exists(pptx_path):
                        self.temp_files.append(pptx_path)  # 記錄臨時文件
                        return pptx_path
        except Exception as e:
            print(f"  LibreOffice 轉換失敗: {e}")
        
        # 方法 2: 在 Windows 上嘗試使用 pywin32
        if os.name == 'nt':
            try:
                import win32com.client
                print(f"  使用 PowerPoint COM 進行轉換...")
                
                powerpoint = win32com.client.Dispatch("PowerPoint.Application")
                powerpoint.Visible = 1
                
                # 打開並轉換
                deck = powerpoint.Presentations.Open(os.path.abspath(ppt_path))
                deck.SaveAs(os.path.abspath(pptx_path), 24)  # 24 = ppSaveAsOpenXMLPresentation
                deck.Close()
                powerpoint.Quit()

                if os.path.exists(pptx_path):
                    self.temp_files.append(pptx_path)  # 記錄臨時文件
                    return pptx_path
            except Exception as e:
                print(f"  COM 轉換失敗: {e}")
        
        return None

    def _cleanup_temp_files(self):
        """清理臨時轉換的文件"""
        import os

        for temp_file in self.temp_files:
            try:
                if os.path.exists(temp_file):
                    os.remove(temp_file)
                    print(f"✓ 已清理臨時文件: {temp_file}")
            except Exception as e:
                print(f"⚠️  清理臨時文件失敗 ({temp_file}): {e}")

        self.temp_files.clear()

    def _encode_image(self, image_path: str) -> str:
        """將圖片編碼為 base64
        
        Args:
            image_path: 圖片路徑
            
        Returns:
            Base64 編碼的圖片字串
        """
        with open(image_path, "rb") as image_file:
            return base64.b64encode(image_file.read()).decode('utf-8')
    
    def _extract_images_from_pdf(self, pdf_path: str) -> List[bytes]:
        """從 PDF 提取圖片
        
        Args:
            pdf_path: PDF 文件路徑
            
        Returns:
            圖片二進制數據列表
        """
        images = []
        try:
            import fitz  # PyMuPDF
            pdf_document = fitz.open(pdf_path)
            
            for page_num in range(len(pdf_document)):
                page = pdf_document[page_num]
                image_list = page.get_images()
                
                for img_index, img in enumerate(image_list):
                    xref = img[0]
                    base_image = pdf_document.extract_image(xref)
                    image_bytes = base_image["image"]
                    images.append(image_bytes)
            
            pdf_document.close()
            return images
        except ImportError:
            print("警告: 需要安裝 PyMuPDF 來提取 PDF 圖片: pip install PyMuPDF --break-system-packages")
            return []
        except Exception as e:
            print(f"提取 PDF 圖片時發生錯誤: {e}")
            return []
    
    def _extract_images_from_docx(self, docx_path: str) -> List[bytes]:
        """從 DOCX 提取圖片
        
        Args:
            docx_path: DOCX 文件路徑
            
        Returns:
            圖片二進制數據列表
        """
        images = []
        try:
            import docx
            from docx.oxml import parse_xml
            
            doc = docx.Document(docx_path)
            
            # 從關係中提取圖片
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    image_bytes = rel.target_part.blob
                    images.append(image_bytes)
            
            return images
        except ImportError:
            print("警告: 需要安裝 python-docx: pip install python-docx --break-system-packages")
            return []
        except Exception as e:
            print(f"提取 DOCX 圖片時發生錯誤: {e}")
            return []
    
    def _extract_images_from_pptx(self, pptx_path: str) -> List[bytes]:
        """從 PPTX 提取圖片
        
        Args:
            pptx_path: PPTX 文件路徑
            
        Returns:
            圖片二進制數據列表
        """
        images = []
        try:
            from pptx import Presentation
            
            prs = Presentation(pptx_path)
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "image"):
                        image_bytes = shape.image.blob
                        images.append(image_bytes)
            
            return images
        except ImportError:
            print("警告: 需要安裝 python-pptx: pip install python-pptx --break-system-packages")
            return []
        except Exception as e:
            print(f"提取 PPTX 圖片時發生錯誤: {e}")
            return []
    
    def read_report(self, file_path: str) -> Tuple[str, List[Dict]]:
        """讀取 FA 報告文件（文字和圖片）
        
        Args:
            file_path: 報告文件路徑
            
        Returns:
            (文字內容, 圖片列表)
        """
        file_path = Path(file_path)
        
        if not file_path.exists():
            raise FileNotFoundError(f"找不到文件: {file_path}")
        
        text_content = ""
        images = []
        
        suffix = file_path.suffix.lower()
        
        # 純圖片文件
        if suffix in ['.jpg', '.jpeg', '.png', '.gif', '.webp']:
            with open(file_path, 'rb') as f:
                image_data = base64.b64encode(f.read()).decode('utf-8')
                images.append({
                    'type': 'image',
                    'data': image_data,
                    'format': suffix[1:]  # 去掉點
                })
            text_content = f"[圖片文件: {file_path.name}]"
        
        # 文字文件
        elif suffix == '.txt':
            with open(file_path, 'r', encoding='utf-8') as f:
                text_content = f.read()
        
        # PDF 文件
        elif suffix == '.pdf':
            try:
                import PyPDF2
                with open(file_path, 'rb') as f:
                    reader = PyPDF2.PdfReader(f)
                    text_content = ""
                    for page in reader.pages:
                        text_content += page.extract_text() + "\n"
                
                # 提取圖片
                image_bytes_list = self._extract_images_from_pdf(str(file_path))
                for img_bytes in image_bytes_list:
                    img_b64 = base64.b64encode(img_bytes).decode('utf-8')
                    images.append({
                        'type': 'image',
                        'data': img_b64,
                        'format': 'png'
                    })
                
            except ImportError:
                print("警告: 需要安裝 PyPDF2: pip install PyPDF2 --break-system-packages")
                raise
        
        # Word 文件
        elif suffix in ['.doc', '.docx']:
            try:
                import docx
                doc = docx.Document(file_path)
                text_content = "\n".join([paragraph.text for paragraph in doc.paragraphs])
                
                # 提取圖片
                image_bytes_list = self._extract_images_from_docx(str(file_path))
                for img_bytes in image_bytes_list:
                    img_b64 = base64.b64encode(img_bytes).decode('utf-8')
                    images.append({
                        'type': 'image',
                        'data': img_b64,
                        'format': 'png'
                    })
                
            except ImportError:
                print("警告: 需要安裝 python-docx: pip install python-docx --break-system-packages")
                raise
        
        # PowerPoint 文件
        elif suffix in ['.ppt', '.pptx']:
            # 處理舊版 .ppt 格式
            if suffix == '.ppt':
                print(f"⚠️  檢測到舊版 PowerPoint 格式 (.ppt)")
                print(f"正在嘗試轉換為 .pptx 格式...")
                
                converted_path = self._convert_ppt_to_pptx(str(file_path))
                if converted_path:
                    print(f"✓ 轉換成功: {converted_path}")
                    file_path = Path(converted_path)
                else:
                    print("\n" + "=" * 70)
                    print("⚠️  無法自動轉換 .ppt 文件")
                    print("=" * 70)
                    print("\n請使用以下方法之一:\n")
                    print("方法 1: 手動轉換 (推薦)")
                    print("  1. 在 PowerPoint 中開啟此文件")
                    print("  2. 另存為 .pptx 格式")
                    print("  3. 重新執行分析\n")
                    print("方法 2: 使用 LibreOffice 轉換")
                    print("  安裝: brew install libreoffice  # macOS")
                    print("       sudo apt install libreoffice  # Linux")
                    print(f"  轉換: libreoffice --headless --convert-to pptx \"{file_path}\"\n")
                    print("方法 3: 線上轉換")
                    print("  使用 CloudConvert 或其他線上轉換工具")
                    print("=" * 70)
                    raise ValueError("不支援 .ppt 格式，請先轉換為 .pptx")
            
            try:
                from pptx import Presentation
                prs = Presentation(file_path)
                
                text_parts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_parts.append(shape.text)
                
                text_content = "\n".join(text_parts)
                
                # 提取圖片
                image_bytes_list = self._extract_images_from_pptx(str(file_path))
                for img_bytes in image_bytes_list:
                    img_b64 = base64.b64encode(img_bytes).decode('utf-8')
                    images.append({
                        'type': 'image',
                        'data': img_b64,
                        'format': 'png'
                    })
                
            except ImportError:
                print("警告: 需要安裝 python-pptx: pip install python-pptx --break-system-packages")
                raise
        
        else:
            raise ValueError(f"不支援的文件格式: {suffix}")
        
        return text_content, images
    
    def create_analysis_prompt(self, report_content: str, has_images: bool = False) -> str:
        """創建分析提示詞
        
        Args:
            report_content: 報告內容
            has_images: 是否包含圖片
            
        Returns:
            完整的分析提示詞
        """
        image_note = ""
        if has_images:
            image_note = """
【注意】此報告包含圖片,請仔細分析圖片中的內容:
- 檢查圖片的清晰度和標註
- 評估圖片是否充分支持分析結論
- 判斷圖表/數據視覺化的品質
"""
        
        prompt = f"""請分析這份 Failure Analysis Report,並根據以下評估維度進行全面評分:
{image_note}
【評估維度與權重】
1. **基本資訊完整性** (15%)
   - 產品資訊(型號、批號、製造日期)
   - 客戶資訊與投訴內容
   - FA 編號與日期
   - 負責工程師資訊

2. **問題描述與定義** (15%)
   - 失效現象描述的清晰度
   - 失效模式的準確性
   - 問題範圍與影響評估
   - 失效率數據

3. **分析方法與流程** (20%)
   - 分析方法的適當性(如:光學檢查、SEM、FIB、X-ray等)
   - 分析步驟的邏輯性與完整性
   - 實驗設計的合理性
   - 分析設備使用的正確性

4. **數據與證據支持** (20%)
   - 分析數據的充分性
   - 圖片/圖表的清晰度與標註
   - 量化數據的準確性
   - 對照組/比較樣本的使用

5. **根因分析** (20%)
   - 根本原因的深度與準確度
   - 因果關係的邏輯推導
   - 5-Why 或 Fishbone 分析的應用
   - 排除其他可能原因的論證

6. **改善對策** (10%)
   - 短期與長期對策的完整性
   - 對策的可行性與有效性
   - 預防措施的提出
   - 驗證計畫

【評分標準】
- **A級 (90-100分)**:卓越報告
- **B級 (80-89分)**:良好報告
- **C級 (70-79分)**:合格報告
- **D級 (60-69分)**:待改進報告
- **F級 (<60分)**:不合格報告

請以 JSON 格式回傳評估結果,格式如下:

{{
  "total_score": <總分數字>,
  "grade": "<等級字母>",
  "dimension_scores": {{
    "基本資訊完整性": {{"score": <分數>, "percentage": <百分比數字>, "comment": "<評語>"}},
    "問題描述與定義": {{"score": <分數>, "percentage": <百分比數字>, "comment": "<評語>"}},
    "分析方法與流程": {{"score": <分數>, "percentage": <百分比數字>, "comment": "<評語>"}},
    "數據與證據支持": {{"score": <分數>, "percentage": <百分比數字>, "comment": "<評語>"}},
    "根因分析": {{"score": <分數>, "percentage": <百分比數字>, "comment": "<評語>"}},
    "改善對策": {{"score": <分數>, "percentage": <百分比數字>, "comment": "<評語>"}}
  }},
  "strengths": [
    "<具體優點1>",
    "<具體優點2>",
    "<具體優點3>"
  ],
  "improvements": [
    {{"priority": "高", "item": "<待改進項目>", "suggestion": "<具體改善建議>"}},
    {{"priority": "中", "item": "<待改進項目>", "suggestion": "<具體改善建議>"}}
  ],
  "summary": "<總評與建議>"
}}

重要格式要求:
1. 你的回應必須是純 JSON 格式,不要包含任何其他文字、markdown 標記或程式碼區塊符號
2. 所有數字欄位(total_score, score, percentage)必須是純數字,不要加單位或符號(例如: 85.5 而不是 85.5% 或 85.5分)
3. percentage 是百分比數值(0-100),例如: 93.33 表示 93.33%
4. 使用台灣繁體中文回答

【FA 報告內容】
{report_content}
"""
        return prompt
    
    def _clean_json_response(self, response_text: str) -> str:
        """清理 AI 返回的 JSON 響應

        Args:
            response_text: 原始響應文本

        Returns:
            清理後的 JSON 字串
        """
        import re

        # 移除 markdown 代碼塊標記
        response_text = response_text.replace('```json', '').replace('```', '').strip()

        # 提取 JSON 部分（處理前後有額外文字的情況）
        if '{' in response_text and '}' in response_text:
            start = response_text.find('{')
            end = response_text.rfind('}') + 1
            response_text = response_text[start:end]

        # 清理 JSON 中的常見格式錯誤
        # 移除數字後面的 % 符號 (例如: "percentage": 93.33% -> "percentage": 93.33)
        response_text = re.sub(r':\s*(\d+\.?\d*)\s*%', r': \1', response_text)

        return response_text

    def analyze_with_ai(self, report_content: str, images: List[Dict] = None) -> Dict:
        """使用 AI 分析報告

        Args:
            report_content: 報告文字內容
            images: 圖片列表

        Returns:
            分析結果字典
        """
        # 根據 skip_images 設定決定是否使用圖片
        if self.skip_images and images:
            print("⚠️  已啟用 --skip-images,將僅分析文字內容")
            images = None

        has_images = images and len(images) > 0
        prompt = self.create_analysis_prompt(report_content, has_images)

        try:
            if self.backend == "ollama":
                return self._analyze_with_ollama(prompt, images)
            elif self.backend == "openai":
                return self._analyze_with_openai(prompt, images)
            elif self.backend == "anthropic":
                return self._analyze_with_anthropic(prompt, images)
            else:
                raise ValueError(f"不支援的後端: {self.backend}")

        except json.JSONDecodeError as e:
            print(f"JSON 解析錯誤: {e}")
            raise
        except Exception as e:
            print(f"分析過程發生錯誤: {e}")
            raise
    
    def _analyze_with_ollama(self, prompt: str, images: List[Dict] = None) -> Dict:
        """使用 Ollama 進行分析"""
        messages = []
        
        # 構建消息內容
        if images and len(images) > 0:
            # 多模態消息
            content_parts = [prompt]
            for img in images[:5]:  # 限制最多 5 張圖片
                content_parts.append({
                    'type': 'image',
                    'data': img['data']
                })
            messages.append({
                'role': 'user',
                'content': prompt,
                'images': [img['data'] for img in images[:5]]
            })
        else:
            messages.append({
                'role': 'user',
                'content': prompt
            })
        
        # 調用 Ollama
        response = ollama.chat(
            model=self.model,
            messages=messages
        )
        
        response_text = response['message']['content'].strip()

        print("=== Ollama raw response ===")
        print(response_text)
        print("=== End raw response ===")

        # 清理並解析 JSON
        response_text = self._clean_json_response(response_text)

        try:
            result = json.loads(response_text)
            return result
        except json.JSONDecodeError as e:
            print(f"\nJSON 解析錯誤: {e}")
            print(f"清理後的響應文本:")
            print(response_text)
            print(f"\n錯誤位置附近的內容:")
            if e.pos < len(response_text):
                start = max(0, e.pos - 50)
                end = min(len(response_text), e.pos + 50)
                print(f"...{response_text[start:end]}...")
                print(f"    {' ' * (min(50, e.pos - start))}^")
            raise
    
    def _analyze_with_openai(self, prompt: str, images: List[Dict] = None) -> Dict:
        """使用 OpenAI API 進行分析"""
        messages = []

        # 構建消息內容
        content = []
        content.append({
            "type": "text",
            "text": prompt
        })

        if images and len(images) > 0:
            for img in images[:10]:  # 限制最多 10 張圖片
                content.append({
                    "type": "image_url",
                    "image_url": {
                        "url": f"data:image/{img['format']};base64,{img['data']}"
                    }
                })

        messages.append({
            "role": "user",
            "content": content
        })

        # 調用 OpenAI
        try:
            response = self.client.chat.completions.create(
                model=self.model,
                messages=messages,
                max_tokens=4000
            )

            response_text = response.choices[0].message.content.strip()

            print("=== OpenAI raw response ===")
            print(response_text)
            print("=== End raw response ===")

            # 檢查是否被拒絕
            if "I'm sorry" in response_text or "I cannot" in response_text or "I can't" in response_text:
                print("\n" + "=" * 80)
                print("⚠️  OpenAI 內容審核拒絕了此請求")
                print("=" * 80)
                print("\n可能原因:")
                print("1. 圖片內容觸發了安全過濾器")
                print("2. 技術術語被誤判為敏感內容")
                print("3. 圖片與文字組合觸發了限制\n")
                print("建議解決方案:")
                print("1. 嘗試不含圖片的純文字分析:")
                print("   python fa_report_analyzer_v2.py -i <文字檔>.txt -b openai -k YOUR_KEY")
                print("\n2. 使用 Ollama 本地模型 (無內容限制):")
                print("   python fa_report_analyzer_v2.py -i <檔案> -b ollama")
                print("\n3. 使用 Anthropic Claude (較少限制):")
                print("   python fa_report_analyzer_v2.py -i <檔案> -b anthropic -k YOUR_KEY")
                print("=" * 80 + "\n")
                raise ValueError("OpenAI API 拒絕處理此請求,請嘗試其他後端或純文字分析")

            # 清理並解析 JSON
            response_text = self._clean_json_response(response_text)

            result = json.loads(response_text)
            return result

        except json.JSONDecodeError as e:
            print("\n" + "=" * 80)
            print("⚠️  OpenAI 返回了無效的 JSON 格式")
            print("=" * 80)
            print(f"JSON 解析錯誤: {e}")
            print(f"清理後的響應文本:")
            print(response_text[:500])
            print(f"\n錯誤位置附近的內容:")
            if e.pos < len(response_text):
                start = max(0, e.pos - 50)
                end = min(len(response_text), e.pos + 50)
                print(f"...{response_text[start:end]}...")
                print(f"    {' ' * (min(50, e.pos - start))}^")
            print("\n這通常表示:")
            print("1. 模型拒絕了請求")
            print("2. 回應格式不符合預期")
            print("\n建議: 嘗試使用其他後端 (ollama 或 anthropic)")
            print("=" * 80 + "\n")
            raise
    
    def _analyze_with_anthropic(self, prompt: str, images: List[Dict] = None) -> Dict:
        """使用 Anthropic Claude 進行分析"""
        content = []
        
        # 添加圖片
        if images and len(images) > 0:
            for img in images[:20]:  # Claude 支持較多圖片
                content.append({
                    "type": "image",
                    "source": {
                        "type": "base64",
                        "media_type": f"image/{img['format']}",
                        "data": img['data']
                    }
                })
        
        # 添加文字提示
        content.append({
            "type": "text",
            "text": prompt
        })
        
        # 調用 Claude
        message = self.client.messages.create(
            model=self.model,
            max_tokens=4000,
            messages=[
                {"role": "user", "content": content}
            ]
        )
        
        response_text = message.content[0].text.strip()

        print("=== Anthropic Claude raw response ===")
        print(response_text)
        print("=== End raw response ===")

        # 清理並解析 JSON
        response_text = self._clean_json_response(response_text)

        try:
            result = json.loads(response_text)
            return result
        except json.JSONDecodeError as e:
            print(f"\nJSON 解析錯誤: {e}")
            print(f"清理後的響應文本:")
            print(response_text)
            print(f"\n錯誤位置附近的內容:")
            if e.pos < len(response_text):
                start = max(0, e.pos - 50)
                end = min(len(response_text), e.pos + 50)
                print(f"...{response_text[start:end]}...")
                print(f"    {' ' * (min(50, e.pos - start))}^")
            raise
    
    def calculate_grade(self, total_score: float) -> Tuple[str, str]:
        """計算等級"""
        for grade, (min_score, max_score, description) in self.grade_criteria.items():
            if min_score <= total_score <= max_score:
                return grade, description
        return 'F', '不合格報告'
    
    def generate_report(self, analysis_result: Dict, output_path: str = None, source_file: str = None) -> str:
        """生成評估報告"""
        total_score = float(analysis_result['total_score'])
        grade = analysis_result['grade']
        grade_desc = [desc for g, (_, _, desc) in self.grade_criteria.items() if g == grade][0]

        # 生成報告內容
        report = []
        report.append("=" * 80)
        report.append("FA 報告評估結果")
        report.append("=" * 80)
        report.append(f"\n生成時間: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
        if source_file:
            report.append(f"來源檔案: {source_file}")
        report.append(f"分析引擎: {self.backend.upper()} ({self.model})\n")
        
        # 1. 總分與等級
        report.append("【總分與等級】")
        report.append(f"總分: {total_score:.1f} 分")
        report.append(f"等級: {grade} 級 - {grade_desc}")
        report.append("")
        
        # 2. 各維度評分表
        report.append("【各維度評分表】")
        report.append("-" * 80)
        
        dimension_data = []
        for dim_name, dim_info in analysis_result['dimension_scores'].items():
            dimension_data.append({
                '評估維度': dim_name,
                '權重': f"{self.dimensions[dim_name]}%",
                '得分': f"{float(dim_info['score']):.1f}",
                '完成度': f"{float(dim_info['percentage']):.1f}%",
                '評語': dim_info['comment']
            })
        
        df = pd.DataFrame(dimension_data)
        report.append(df.to_string(index=False))
        report.append("")
        
        # 3. 優點分析
        report.append("【優點分析】")
        for i, strength in enumerate(analysis_result['strengths'], 1):
            report.append(f"{i}. {strength}")
        report.append("")
        
        # 4. 待改進項目
        report.append("【待改進項目】(依優先級排序)")
        for i, item in enumerate(analysis_result['improvements'], 1):
            report.append(f"\n{i}. [{item['priority']}優先級] {item['item']}")
            report.append(f"   改善建議: {item['suggestion']}")
        report.append("")
        
        # 5. 總評與建議
        report.append("【總評與建議】")
        report.append(analysis_result['summary'])
        report.append("")
        
        report.append("=" * 80)
        report.append("報告結束")
        report.append("=" * 80)
        
        report_text = "\n".join(report)
        
        # 保存到文件
        if output_path:
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(report_text)
            print(f"\n✓ 評估報告已保存至: {output_path}")
        
        return report_text
    
    def analyze_report(self, input_file: str, output_file: str = None) -> Dict:
        """完整的報告分析流程"""
        import os

        print("=" * 80)
        print("FA 報告分析工具 v2.0")
        print("=" * 80)

        try:
            # 1. 讀取報告
            print(f"\n[1/3] 讀取報告文件: {input_file}")
            report_content, images = self.read_report(input_file)
            print(f"✓ 成功讀取報告 ({len(report_content)} 字元)")
            if images:
                print(f"✓ 提取到 {len(images)} 張圖片")

            # 2. AI 分析
            print(f"\n[2/3] 使用 {self.backend.upper()} 進行深度分析...")
            analysis_result = self.analyze_with_ai(report_content, images)
            print("✓ 分析完成")

            # 3. 生成報告
            print("\n[3/3] 生成評估報告...")

            # 創建輸出資料夾
            output_dir = "evaluation_results"
            os.makedirs(output_dir, exist_ok=True)

            # 如果沒有指定輸出文件，自動生成文件名
            if not output_file:
                output_file = f"fa_evaluation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.txt"

            # 確保輸出文件在指定資料夾內
            if not output_file.startswith(output_dir):
                output_file = os.path.join(output_dir, os.path.basename(output_file))

            # 獲取來源文件名
            source_filename = os.path.basename(input_file)

            report_text = self.generate_report(analysis_result, output_file, source_filename)

            print("\n" + "=" * 80)
            print("分析完成!")
            print("=" * 80)

            return analysis_result

        finally:
            # 清理臨時轉換的文件（無論分析成功或失敗都會執行）
            if self.temp_files:
                print("\n[清理] 移除臨時轉換文件...")
                self._cleanup_temp_files()


def main():
    """主程式"""
    import argparse
    
    parser = argparse.ArgumentParser(
        description='FA Report Analyzer v2.0 - 支援 Ollama/OpenAI/Anthropic 和圖片解析',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
範例用法:
  # 使用 Ollama (預設)
  python fa_report_analyzer_v2.py -i fa_report.pdf
  
  # 使用 OpenAI API
  python fa_report_analyzer_v2.py -i report.pdf -b openai -k YOUR_API_KEY
  
  # 使用 Anthropic Claude
  python fa_report_analyzer_v2.py -i report.pdf -b anthropic -k YOUR_API_KEY
  
  # 指定 Ollama 模型
  python fa_report_analyzer_v2.py -i report.pdf -b ollama -m llama3.2-vision
  
  # 分析包含圖片的報告
  python fa_report_analyzer_v2.py -i report_with_images.pdf -o evaluation.txt
        """
    )
    
    parser.add_argument('-i', '--input', required=True,
                        help='輸入的 FA 報告文件路徑')
    parser.add_argument('-o', '--output',
                        help='輸出的評估報告文件路徑 (預設: 自動生成)')
    parser.add_argument('-b', '--backend', default='ollama',
                        choices=['ollama', 'openai', 'anthropic'],
                        help='LLM 後端 (預設: ollama)')
    parser.add_argument('-m', '--model',
                        help='模型名稱 (預設: 依後端自動選擇)')
    parser.add_argument('-k', '--api-key',
                        help='API key (OpenAI/Anthropic 需要)')
    parser.add_argument('--base-url',
                        help='API base URL (OpenAI 相容接口)')
    parser.add_argument('--skip-images', action='store_true',
                        help='跳過圖片分析,僅分析文字內容 (可避免 OpenAI 內容審核問題)')

    args = parser.parse_args()

    try:
        # 創建分析器
        analyzer = FAReportAnalyzer(
            backend=args.backend,
            model=args.model,
            api_key=args.api_key,
            base_url=args.base_url,
            skip_images=args.skip_images
        )
        
        # 執行分析
        result = analyzer.analyze_report(args.input, args.output)
        
        # 顯示摘要
        print(f"\n總分: {float(result['total_score']):.1f} 分")
        print(f"等級: {result['grade']}")
        
    except Exception as e:
        print(f"\n錯誤: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == "__main__":
    main()
